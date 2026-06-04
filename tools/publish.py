#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
============================================================
 publish.py — 一键发布脚本
============================================================
把一份 PPT 一条命令发布到博客「数据观察」+ 生成公众号素材。

自动完成：
  ① 从 PPT 自动提取元信息（标题/副标题/摘要/标签/页数）
  ② 复制 PPT 到 blog/decks/YYYY-MM-DD.pptx
  ③ 更新 data/decks.js（自动插入记录，id 自增）
  ④ 生成公众号素材（图片 + 文案）
  ⑤ git 提交并尝试推送上线

用法：
  python tools/publish.py 数字消费研究日报_20260604.pptx
  python tools/publish.py "D:\\Claude code\\自动规划\\xxx.pptx"

可选参数：
  --dry-run       只提取并打印元信息，不改任何文件（先看效果用）
  --skip-wechat   跳过公众号素材生成
  --no-push       只本地提交，不推送（之后自己用 GitHub Desktop 推）
  --no-wx-draft   生成公众号素材但不自动发草稿（留着手动发）
============================================================
"""
import os
import re
import sys
import shutil
import subprocess
from pathlib import Path

# Windows 控制台 UTF-8 输出
try:
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")
except Exception:
    pass

from pptx import Presentation

# ===== 路径常量 =====
BLOG_DIR = Path(__file__).resolve().parent.parent     # blog/
DECKS_DIR = BLOG_DIR / "decks"                         # blog/decks/
DATA_JS = BLOG_DIR / "data" / "decks.js"               # blog/data/decks.js
DEFAULT_SRC = Path(r"D:\Claude code\自动规划")          # PPT 默认来源目录

# ===== 标签关键词库（命中即作为标签）=====
TAG_KEYWORDS = [
    "即时零售", "跨境电商", "直播电商", "社区团购", "内容消费",
    "AI购物", "会展经济", "银发经济", "下沉市场", "绿色消费",
    "政策评估", "空间计量", "因果识别", "消费趋势", "数字金融",
]
BASE_TAGS = ["数字消费", "研究日报"]


def parse_date(name: str) -> str:
    """从文件名里抓 8 位日期：20260604 → 2026-06-04。"""
    m = re.search(r"(20\d{6})", name)
    if not m:
        raise ValueError(f"文件名里找不到 8 位日期（如 20260604）：{name}")
    s = m.group(1)
    return f"{s[:4]}-{s[4:6]}-{s[6:8]}"


def good_line(s: str) -> bool:
    """判断一行文字是否适合作为摘要要点：长度适中、不是栏目标题/来源标注。"""
    s = s.strip()
    if not (14 <= len(s) <= 64):
        return False
    c = s[0]
    # 排除 emoji / 符号开头的栏目标题；允许中文或数字开头
    if not (c.isdigit() or "一" <= c <= "鿿"):
        return False
    bad = ["来源", "http", "📍", "📅", "声明", "Andy", "数据说明", "研究日报  ·", "转载"]
    if any(b in s for b in bad):
        return False
    # 含完整日期（YYYY-MM-DD）的行通常是来源/时间标注，不作为摘要要点
    if re.search(r"\d{4}-\d{2}-\d{2}", s):
        return False
    return True


def extract_meta(pptx_path: Path, date_tag: str) -> dict:
    """从 PPT 自动提取元信息。"""
    prs = Presentation(str(pptx_path))
    slides = []
    for slide in prs.slides:
        texts = [sh.text_frame.text.strip() for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        slides.append(texts)

    pages = len(slides)
    y, m, d = date_tag[:4], date_tag[5:7], date_tag[8:10]
    title = f"数字消费研究日报 · {y}/{m}/{d}"

    # 副标题：第一页里含「专题」的句子，取冒号后内容
    subtitle = ""
    first = slides[0] if slides else []
    for t in first:
        if "专题" in t:
            subtitle = t.split("：", 1)[-1].strip() if "：" in t else t.strip()
            break
    if not subtitle:
        subtitle = "每日数字消费市场观察"

    # 摘要：主题句 + 前 3 条要点
    highlights, seen = [], set()
    for texts in slides[1:]:           # 跳过封面页
        for t in texts:
            t = t.strip()
            if good_line(t) and t not in seen:
                seen.add(t)
                highlights.append(t)
    top = highlights[:3]
    if top:
        summary = f"今日聚焦「{subtitle}」。" + "；".join(top) + "。"
    else:
        summary = f"今日聚焦「{subtitle}」。"
    summary = summary[:180]

    # 标签：基础标签 + 命中关键词
    full = " ".join(" ".join(ts) for ts in slides)
    hits = [k for k in TAG_KEYWORDS if k in full]
    tags = BASE_TAGS + [k for k in hits if k not in BASE_TAGS]
    tags = tags[:5]

    return {
        "title": title, "subtitle": subtitle, "summary": summary,
        "tags": tags, "pages": pages,
    }


def js_escape(s: str) -> str:
    """转义 JS 字符串里的双引号和反斜杠。"""
    return s.replace("\\", "\\\\").replace('"', '\\"')


def update_decks_js(meta: dict, date_tag: str, file_name: str) -> int:
    """在 decks.js 的 DECKS 数组最前面插入一条记录，返回新 id。"""
    content = DATA_JS.read_text(encoding="utf-8")

    if f'"{file_name}"' in content:
        ids = [int(x) for x in re.findall(r"id:\s*(\d+)", content)]
        print(f"⚠ decks.js 里已存在 {file_name} 的记录，跳过插入。")
        return max(ids) if ids else 1

    ids = [int(x) for x in re.findall(r"id:\s*(\d+)", content)]
    new_id = (max(ids) + 1) if ids else 1
    tags_str = ", ".join(f'"{js_escape(t)}"' for t in meta["tags"])

    record = (
        "  {\n"
        f"    id: {new_id},\n"
        f'    date: "{date_tag}",\n'
        f'    title: "{js_escape(meta["title"])}",\n'
        f'    subtitle: "{js_escape(meta["subtitle"])}",\n'
        f'    summary: "{js_escape(meta["summary"])}",\n'
        f"    tags: [{tags_str}],\n"
        f"    pages: {meta['pages']},\n"
        f'    file: "{file_name}",\n'
        '    pdf: ""\n'
        "  },\n"
    )

    # 在 "const DECKS = [\n" 之后插入新记录（用 lambda 避免替换串转义问题）
    new_content, n = re.subn(
        r"(const DECKS = \[\n)",
        lambda mt: mt.group(1) + record,
        content, count=1,
    )
    if n == 0:
        raise RuntimeError("没找到 `const DECKS = [`，decks.js 结构可能被改过。")
    DATA_JS.write_text(new_content, encoding="utf-8")
    return new_id


def run_wechat(date_tag: str):
    """调用 wechat-cards.py 生成公众号素材（图片 + 文案）。"""
    script = BLOG_DIR / "tools" / "wechat-cards.py"
    subprocess.run([sys.executable, str(script), f"{date_tag}.pptx"],
                   cwd=str(BLOG_DIR))


def run_wechat_draft(date_tag: str):
    """调用 auto-publish.ts，把素材发到公众号草稿箱（需 bun + 已配置公众号凭据）。"""
    script = BLOG_DIR / "tools" / "auto-publish.ts"
    if not script.exists():
        print("⚠ 未找到 tools/auto-publish.ts，跳过公众号草稿发布。")
        return
    bun = shutil.which("bun")
    if not bun:
        print("⚠ 未安装 bun，跳过公众号草稿发布。素材已生成，可手动发。")
        return
    print("\n📤 发布到公众号草稿箱...")
    r = subprocess.run([bun, str(script), date_tag], cwd=str(BLOG_DIR))
    if r.returncode == 0:
        print("✅ 已发到公众号草稿箱！登录公众号后台点「发送」即可。")
    else:
        print("⚠ 公众号草稿发布失败。素材已在 tools/wechat-output/，可手动发布。")


def git_publish(date_tag: str, push: bool):
    """git add + commit (+ push)。push 失败不报错，提示改用 GitHub Desktop。"""
    def run(args):
        return subprocess.run(["git"] + args, cwd=str(BLOG_DIR),
                              capture_output=True, text=True)

    run(["add", "-A"])
    c = run(["commit", "-m", f"自动发布：{date_tag} 数字消费日报"])
    if c.returncode != 0 and "nothing to commit" in (c.stdout + c.stderr):
        print("ℹ 没有需要提交的改动。")
        return
    print("✅ 已本地提交。")

    if not push:
        print("⏭  已跳过推送（--no-push）。请用 GitHub Desktop 点 Push 上线。")
        return

    p = run(["push"])
    if p.returncode == 0:
        print("🚀 已推送上线！1-2 分钟后访问 https://w2950260544-ai.github.io/decks.html")
    else:
        print("⚠ 自动推送失败（常见于网络/认证）。请打开 GitHub Desktop 点 Push 即可。")
        err = (p.stderr or "").strip().splitlines()
        if err:
            print("   详情：", err[-1][:160])


def resolve_pptx(arg: str) -> Path:
    """把参数解析成 PPT 完整路径：支持完整路径或仅文件名（去默认目录找）。"""
    p = Path(arg)
    if p.is_file():
        return p
    cand = DEFAULT_SRC / arg
    if cand.is_file():
        return cand
    if not arg.lower().endswith(".pptx"):
        cand2 = DEFAULT_SRC / (arg + ".pptx")
        if cand2.is_file():
            return cand2
    raise FileNotFoundError(f"找不到 PPT：{arg}（已在当前目录和 {DEFAULT_SRC} 查找）")


def main():
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    flags = {a for a in sys.argv[1:] if a.startswith("--")}
    if not args:
        # 无参数：自动选默认目录里最新的 PPT（方便双击 bat 运行）
        cands = [p for p in DEFAULT_SRC.glob("*.pptx") if not p.name.startswith("~$")]
        if cands:
            latest = max(cands, key=lambda p: p.stat().st_mtime)
            print(f"ℹ 未指定文件，自动选择最新 PPT：{latest.name}\n")
            args = [str(latest)]
        else:
            print(__doc__)
            print(f"\n（默认目录 {DEFAULT_SRC} 里也没找到 .pptx）")
            sys.exit(1)

    dry_run = "--dry-run" in flags
    skip_wechat = "--skip-wechat" in flags
    no_wx_draft = "--no-wx-draft" in flags
    push = "--no-push" not in flags

    src = resolve_pptx(args[0])
    date_tag = parse_date(src.name)

    print(f"📂 源 PPT：{src}")
    print(f"📅 日期：{date_tag}")
    print("🔍 正在提取元信息...\n")

    meta = extract_meta(src, date_tag)
    print("─" * 50)
    print(f"  标题   : {meta['title']}")
    print(f"  副标题 : {meta['subtitle']}")
    print(f"  页数   : {meta['pages']}")
    print(f"  标签   : {' / '.join(meta['tags'])}")
    print(f"  摘要   : {meta['summary']}")
    print("─" * 50 + "\n")

    if dry_run:
        print("✋ --dry-run 模式：仅预览元信息，未改动任何文件。")
        print("   满意的话去掉 --dry-run 重新运行即可正式发布。")
        return

    # ② 复制 PPT
    DECKS_DIR.mkdir(parents=True, exist_ok=True)
    dest = DECKS_DIR / f"{date_tag}.pptx"
    shutil.copy(src, dest)
    print(f"✅ 已复制 PPT → decks/{dest.name}")

    # ③ 更新 decks.js
    new_id = update_decks_js(meta, date_tag, dest.name)
    print(f"✅ 已更新 decks.js（id={new_id}）")

    # ④ 公众号素材
    if skip_wechat:
        print("⏭  已跳过公众号素材生成（--skip-wechat）。")
    else:
        print("\n📱 生成公众号素材...")
        run_wechat(date_tag)

    # ⑤ git（博客上线）
    print("\n📤 提交到 Git...")
    git_publish(date_tag, push)

    # ⑥ 公众号草稿箱（接力 auto-publish.ts）
    if not skip_wechat and not no_wx_draft:
        run_wechat_draft(date_tag)

    print("\n🎉 全部完成！")
    print(f"   博客链接：https://w2950260544-ai.github.io/deck.html?id={new_id}")
    if not skip_wechat:
        print(f"   公众号素材：tools/wechat-output/{date_tag}/")
        if not no_wx_draft:
            print("   公众号草稿：已尝试发到草稿箱，去后台点「发送」即可。")


if __name__ == "__main__":
    main()
